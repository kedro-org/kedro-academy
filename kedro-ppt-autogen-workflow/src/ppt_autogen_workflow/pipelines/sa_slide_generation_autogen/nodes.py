"""Kedro nodes for single-agent AutoGen PPT generation pipeline."""
from __future__ import annotations

import asyncio
import logging
import tempfile
from typing import Any
from pathlib import Path

import matplotlib
matplotlib.use('Agg')

import pandas as pd
import matplotlib.pyplot as plt

from autogen_ext.models.openai import OpenAIChatCompletionClient
from pptx import Presentation

from .agent import create_ppt_agent, PPTGenerationAgent
from .tools import build_tools
from ppt_autogen_workflow.utils.ppt_builder import create_slide, combine_presentations
from ppt_autogen_workflow.utils.instruction_parser import parse_instructions_yaml
from ppt_autogen_workflow.utils.node_helpers import (
    format_summary_text,
    extract_chart_path,
    extract_summary_text,
    create_fallback_summary,
)

logger = logging.getLogger(__name__)


def init_tools(sales_data: pd.DataFrame) -> dict[str, Any]:
    """Initialize all tools needed for PPT generation agent."""
    return {"ppt_tools": build_tools(sales_data)}


def compile_ppt_agent(
    slide_generation_requirements: dict[str, Any],
    ppt_generator_system_prompt: Any,
    ppt_generator_user_prompt: Any,
    model_client: OpenAIChatCompletionClient,
    tools: dict[str, Any],
) -> PPTGenerationAgent:
    """Compile the PPT Generation Agent with requirements and formatted prompts."""
    try:
        slide_definitions = parse_instructions_yaml(slide_generation_requirements)
        ppt_requirement = {'slides': {}}
        data_context = "Sales data available through agent tools"

        for slide_key, slide_config in slide_definitions.items():
            objective = slide_config.get('objective', {})
            ppt_requirement['slides'][slide_key] = {
                'slide_title': objective.get('slide_title', slide_key),
                'chart_instruction': objective.get('chart_instruction', ''),
                'summary_instruction': objective.get('summary_instruction', ''),
                'data_context': data_context,
            }

        system_prompt_text = ppt_generator_system_prompt.format()
        agent = create_ppt_agent(model_client, system_prompt_text, tools["ppt_tools"])
        agent._ppt_requirement = ppt_requirement

        # Format user prompts for each slide
        formatted_prompts = {}
        for slide_key, slide_config in ppt_requirement['slides'].items():
            slide_config_str = (
                f"Slide Title: {slide_config['slide_title']}\n"
                f"Chart Instruction: {slide_config['chart_instruction']}\n"
                f"Summary Instruction: {slide_config['summary_instruction']}\n"
                f"Data Context: {slide_config['data_context']}"
            )
            formatted_prompt = ppt_generator_user_prompt.format(slide_config=slide_config_str)
            formatted_prompts[slide_key] = formatted_prompt

        agent._formatted_prompts = formatted_prompts
        return agent

    except Exception as e:
        logger.error(f"Failed to compile PPT agent: {str(e)}", exc_info=True)
        raise


def generate_presentation(compiled_ppt_agent: PPTGenerationAgent) -> Any:
    """Generate presentation using compiled agent."""
    try:
        slide_configs = compiled_ppt_agent._ppt_requirement['slides']
        slide_presentations = []

        for slide_key, config in slide_configs.items():
            slide_title = config['slide_title']
            chart_path, _ = _generate_chart(
                compiled_ppt_agent, config.get('chart_instruction', ''), slide_key
            )
            summary_text = _generate_summary(
                compiled_ppt_agent, config.get('summary_instruction', ''), slide_key, chart_path
            )
            formatted_summary = format_summary_text(summary_text)
            slide_prs = _create_slide_presentation(slide_title, chart_path, formatted_summary)
            slide_presentations.append(slide_prs)

        return combine_presentations(slide_presentations)

    except Exception as e:
        logger.error(f"PPT generation failed: {str(e)}", exc_info=True)
        return _create_error_presentation(str(e))


def _generate_chart(
    agent: PPTGenerationAgent, chart_instruction: str, slide_key: str
) -> tuple[str, plt.Figure | None]:
    """Generate chart for a slide using agent."""
    try:
        query = f"Please generate a chart using the generate_sales_chart tool.\n\nChart Instruction: {chart_instruction}"
        chart_result = asyncio.run(agent.invoke(query))
        chart_path = extract_chart_path(chart_result)

        if chart_path and Path(chart_path).exists():
            try:
                from matplotlib.image import imread
                img = imread(chart_path)
                fig, ax = plt.subplots(figsize=(10, 6))
                ax.imshow(img)
                ax.axis('off')
                return str(chart_path), fig
            except Exception as e:
                logger.warning(f"Could not load chart image: {e}")

        fig, ax = plt.subplots(figsize=(10, 6))
        ax.text(0.5, 0.5, f"Chart for {slide_key}\n{chart_instruction[:50]}...",
                ha='center', va='center', fontsize=12, wrap=True)
        ax.set_title(f"Chart: {slide_key}")

        temp_dir = Path(tempfile.mkdtemp())
        chart_path = temp_dir / f"chart_{slide_key}.png"
        fig.savefig(chart_path, dpi=300, bbox_inches='tight')
        return str(chart_path), fig

    except Exception as e:
        logger.error(f"Error generating chart for {slide_key}: {str(e)}")
        return "", None


def _generate_summary(
    agent: PPTGenerationAgent, summary_instruction: str, slide_key: str, chart_path: str = None
) -> str:
    """Generate summary text for a slide using agent."""
    try:
        chart_status = f"Chart generated: {chart_path}" if chart_path and Path(chart_path).exists() else "Chart generation in progress"
        query = f"""Please generate a summary using the generate_business_summary tool.

Summary Instruction: {summary_instruction}
Chart Status: {chart_status}

Use actual calculated values, not placeholders."""

        summary_result = asyncio.run(agent.invoke(query))
        summary_text = extract_summary_text(summary_result)
        return summary_text if summary_text else create_fallback_summary(slide_key, summary_instruction)

    except Exception as e:
        logger.error(f"Error generating summary for {slide_key}: {str(e)}")
        return create_fallback_summary(slide_key, summary_instruction)


def _create_slide_presentation(title: str, chart_path: str | None, summary: str) -> Presentation:
    """Create a single slide presentation."""
    chart = chart_path if chart_path and Path(chart_path).exists() else ""
    return create_slide(slide_title=title, chart_path=chart, summary_text=summary)


def _create_error_presentation(error_message: str) -> Presentation:
    """Create error presentation when generation fails."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Error in Presentation Generation"
    slide.placeholders[1].text = f"Error: {error_message}"
    return prs
