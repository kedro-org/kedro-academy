"""Agent helper functions for single-agent PPT generation pipeline.

This module contains helper functions that work with the PPT generation agent
to generate charts and summaries. These are separated from nodes.py to keep
node functions focused on Kedro pipeline logic while agent-specific operations
are encapsulated here.
"""
from __future__ import annotations

import asyncio
import logging
from pathlib import Path

from pptx import Presentation

from .agent import PPTGenerationAgent
from ppt_autogen_workflow.base import ChartOutput, SummaryOutput
from ppt_autogen_workflow.pipelines.ma_slide_generation_autogen.chart import (
    ChartGenerationError,
)

logger = logging.getLogger(__name__)


def generate_chart(
    agent: PPTGenerationAgent, user_prompt: str, slide_key: str, chart_instruction: str = ""
) -> str:
    """Generate chart for a slide using agent with structured output.

    Args:
        agent: Compiled PPT generation agent
        user_prompt: Formatted user prompt for the slide
        slide_key: Unique identifier for the slide
        chart_instruction: Chart instruction (for error context)

    Returns:
        Path to the generated chart image

    Raises:
        ChartGenerationError: If chart generation fails or produces invalid output
    """
    try:
        query = f"{user_prompt}\n\nTask: Generate a chart using the generate_sales_chart tool."
        chart_output: ChartOutput = asyncio.run(agent.invoke_for_chart(query))
        chart_path = chart_output.chart_path

        if chart_path and Path(chart_path).exists():
            return chart_path

        # Chart generation failed - raise error instead of silent fallback
        raise ChartGenerationError(
            f"Chart generation for '{slide_key}' did not produce a valid file. "
            f"Instruction: {chart_instruction[:100]}..."
        )

    except ChartGenerationError:
        raise
    except Exception as e:
        logger.error(f"Error generating chart for {slide_key}: {str(e)}")
        raise ChartGenerationError(
            f"Chart generation failed for '{slide_key}': {str(e)}"
        ) from e


def generate_summary(
    agent: PPTGenerationAgent, user_prompt: str, slide_key: str, chart_path: str = None
) -> str:
    """Generate summary text for a slide using agent with structured output.

    Args:
        agent: Compiled PPT generation agent
        user_prompt: Formatted user prompt for the slide
        slide_key: Unique identifier for the slide
        chart_path: Path to generated chart for status

    Returns:
        Generated summary text
    """
    try:
        chart_status = f"Chart generated: {chart_path}" if chart_path and Path(chart_path).exists() else "Chart generation in progress"
        query = f"""{user_prompt}

Task: Generate a summary using the generate_business_summary tool.
Chart Status: {chart_status}

Use actual calculated values, not placeholders."""

        summary_output: SummaryOutput = asyncio.run(agent.invoke_for_summary(query))
        summary_text = summary_output.summary_text

        if summary_text:
            return summary_text

        # Fallback
        logger.warning(f"Summary not generated for {slide_key}, using fallback")
        return _create_fallback_summary(slide_key, user_prompt)

    except Exception as e:
        logger.error(f"Error generating summary for {slide_key}: {str(e)}")
        return _create_fallback_summary(slide_key, user_prompt)


def _create_fallback_summary(slide_key: str, instruction: str = "") -> str:
    """Create fallback summary when agent fails."""
    if instruction:
        return f"• Analysis for {slide_key}\n• Data insights based on: {instruction[:100]}..."
    return f"• Analysis for {slide_key}\n• Data insights generated"


def create_error_presentation(error_message: str) -> Presentation:
    """Create error presentation when generation fails.

    Args:
        error_message: Error message to display

    Returns:
        Presentation object with error slide
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Error in Presentation Generation"
    slide.placeholders[1].text = f"Error: {error_message}"
    return prs
