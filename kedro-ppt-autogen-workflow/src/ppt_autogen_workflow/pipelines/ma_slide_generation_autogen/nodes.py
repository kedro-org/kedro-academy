"""Kedro nodes for multi-agent AutoGen PPT generation pipeline.

This module contains node functions with clear separation between:
1. parse_requirements - Deterministic data preparation
2. orchestrate_agents - Agentic chart/summary generation
3. assemble_presentation - Deterministic slide assembly
"""
from __future__ import annotations

import asyncio
import logging
from pathlib import Path
from typing import Any

from pptx import Presentation
from kedro.pipeline.llm_context import LLMContext

# Import from new module structure
from .planner import create_planner_agent
from .chart import create_chart_generator_agent, generate_chart
from .summary import create_summarizer_agent, generate_summary
from .critic import create_critic_agent, run_qa_review
from .presentation import create_slide, combine_presentations, format_summary_text
from .orchestration_helpers import (
    create_agent_from_context,
    format_all_prompts,
)
from ppt_autogen_workflow.utils.instruction_parser import parse_instructions_yaml

logger = logging.getLogger(__name__)


def parse_requirements(
    slide_generation_requirements: dict[str, Any],
) -> dict[str, Any]:
    """Parse slide generation requirements from YAML.

    This is a deterministic node that prepares slide configurations
    for the agents to process.

    Args:
        slide_generation_requirements: Raw slide configuration from YAML

    Returns:
        Parsed slide configurations with agent-specific views
    """
    slide_definitions = parse_instructions_yaml(slide_generation_requirements)
    data_context = "Sales data available through agent tools"

    planner_slides = {}
    chart_slides = {}
    summarizer_slides = {}

    for slide_key, slide_config in slide_definitions.items():
        objective = slide_config.get('objective', {})
        slide_title = objective.get('slide_title', slide_key)
        chart_instruction = objective.get('chart_instruction', '')
        summary_instruction = objective.get('summary_instruction', '')

        planner_slides[slide_key] = {
            'slide_title': slide_title,
            'chart_instruction': chart_instruction,
            'summary_instruction': summary_instruction,
            'data_context': data_context,
        }

        chart_slides[slide_key] = {
            'slide_title': slide_title,
            'chart_instruction': chart_instruction,
            'data_context': data_context,
        }

        summarizer_slides[slide_key] = {
            'slide_title': slide_title,
            'summary_instruction': summary_instruction,
            'data_context': data_context,
        }

    return {
        'planner_slides': planner_slides,
        'chart_slides': chart_slides,
        'summarizer_slides': summarizer_slides,
    }


def orchestrate_multi_agent_workflow(
    planner_context: LLMContext,
    chart_context: LLMContext,
    summarizer_context: LLMContext,
    critic_context: LLMContext,
    slide_configs: dict[str, Any],
    styling_params: dict[str, Any],
    layout_params: dict[str, Any],
    quality_assurance_params: dict[str, Any],
) -> tuple[dict[str, str], dict[str, str], dict[str, Any]]:
    """Orchestrate multi-agent workflow to generate charts and summaries.

    This is the agentic node that coordinates multiple agents to generate content.
    It receives pre-parsed requirements and outputs raw results.

    Args:
        planner_context: LLMContext for planner agent
        chart_context: LLMContext for chart generator agent
        summarizer_context: LLMContext for summarizer agent
        critic_context: LLMContext for critic agent
        slide_configs: Parsed slide configurations from parse_requirements
        styling_params: Styling parameters for presentation
        layout_params: Layout parameters for presentation
        quality_assurance_params: QA parameters for critic

    Returns:
        Tuple of (slide_chart_paths, slide_summaries, slide_configs)
    """
    # Extract parsed requirements
    planner_slides = slide_configs.get('planner_slides', {})
    chart_slides = slide_configs.get('chart_slides', {})
    summarizer_slides = slide_configs.get('summarizer_slides', {})

    # Create agents from LLMContext objects
    planner_agent = create_agent_from_context(
        planner_context, "planner_system_prompt", create_planner_agent
    )
    chart_agent = create_agent_from_context(
        chart_context, "chart_generator_system_prompt", create_chart_generator_agent
    )
    summarizer_agent = create_agent_from_context(
        summarizer_context, "summarizer_system_prompt", create_summarizer_agent
    )
    critic_agent = create_agent_from_context(
        critic_context, "critic_system_prompt", create_critic_agent
    )

    # Format prompts
    critic_user_prompt = critic_context.prompts.get("critic_user_prompt")
    planner_prompts, chart_prompts, summarizer_prompts = format_all_prompts(
        planner_slides, chart_slides, summarizer_slides,
        planner_context.prompts.get("planner_user_prompt"),
        chart_context.prompts.get("chart_generator_user_prompt"),
        summarizer_context.prompts.get("summarizer_user_prompt"),
    )

    # Run multi-agent workflow
    slide_chart_paths = {}
    slide_summaries = {}

    for slide_key, config in planner_slides.items():
        # Planner analyzes requirements
        asyncio.run(planner_agent.invoke(planner_prompts[slide_key]))

        # Chart generation
        chart_path = generate_chart(
            chart_agent, chart_prompts[slide_key], slide_key, config.get('chart_instruction', '')
        )
        slide_chart_paths[slide_key] = chart_path

        # Summary generation
        chart_status = f"Chart generated: {chart_path}" if chart_path else "Chart generation in progress"
        summary_query = summarizer_prompts[slide_key].replace('{chart_status}', chart_status)
        summary_text = generate_summary(
            summarizer_agent, summary_query, slide_key, config.get('summary_instruction', '')
        )
        slide_summaries[slide_key] = format_summary_text(summary_text)

        # QA review
        run_qa_review(
            critic_agent, critic_user_prompt, quality_assurance_params,
            config['slide_title'], chart_path, summary_text, config
        )

    # Package results - include planner_slides for assemble_presentation
    return slide_chart_paths, slide_summaries, {
        'slides': planner_slides,
        'layout': layout_params,
        'styling': styling_params,
    }


def assemble_presentation(
    slide_chart_paths: dict[str, str],
    slide_summaries: dict[str, str],
    slide_configs: dict[str, Any],
) -> Any:
    """Assemble final presentation from generated charts and summaries.

    This is a deterministic node that creates slides from the agent-generated content.

    Args:
        slide_chart_paths: Dict mapping slide_key to chart image path
        slide_summaries: Dict mapping slide_key to formatted summary text
        slide_configs: Dict containing slide metadata and layout/styling params

    Returns:
        Combined PowerPoint presentation
    """
    try:
        slides = slide_configs.get('slides', {})
        layout_params = slide_configs.get('layout', {})
        styling_params = slide_configs.get('styling', {})

        slide_presentations = []
        for slide_key, config in slides.items():
            slide_title = config['slide_title']
            chart_path = slide_chart_paths.get(slide_key, '')
            summary = slide_summaries.get(slide_key, '')

            chart = chart_path if chart_path and Path(chart_path).exists() else ""
            slide_prs = create_slide(
                slide_title=slide_title,
                chart_path=chart,
                summary_text=summary,
                layout_params=layout_params,
                styling_params=styling_params,
            )
            slide_presentations.append(slide_prs)

        return combine_presentations(slide_presentations)

    except Exception as e:
        logger.error(f"Presentation assembly failed: {str(e)}", exc_info=True)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Presentation Assembly Error"
        slide.placeholders[1].text = f"Error: {str(e)}"
        return prs
