"""Agent helper functions for multi-agent PPT generation pipeline.

This module contains helper functions that work with specific agents
to generate charts, summaries, and perform QA reviews. These are
separated from nodes.py to keep node functions focused on Kedro pipeline
logic while agent-specific operations are encapsulated here.
"""
from __future__ import annotations

import asyncio
import logging
import tempfile
from pathlib import Path
from typing import Any

import matplotlib
matplotlib.use('Agg')

import matplotlib.pyplot as plt
from pptx import Presentation

from .agent import ChartGeneratorAgent, SummarizerAgent, CriticAgent
from ppt_autogen_workflow.utils.ppt_builder import create_slide
from ppt_autogen_workflow.utils.node_helpers import (
    extract_chart_path,
    extract_summary_text,
    create_fallback_summary,
)

logger = logging.getLogger(__name__)


def generate_chart(
    chart_agent: ChartGeneratorAgent, query: str, slide_key: str, instruction: str
) -> tuple[str, plt.Figure | None]:
    """Generate chart using chart agent.

    Args:
        chart_agent: Compiled chart generator agent
        query: Formatted query for chart generation
        slide_key: Unique identifier for the slide
        instruction: Chart instruction for fallback

    Returns:
        Tuple of (chart_path, matplotlib_figure)
    """
    try:
        chart_result = asyncio.run(chart_agent.invoke(query))
        chart_path = extract_chart_path(chart_result)

        if chart_path and Path(chart_path).exists():
            try:
                from matplotlib.image import imread
                img = imread(chart_path)
                fig, ax = plt.subplots(figsize=(10, 6))
                ax.imshow(img)
                ax.axis('off')
                return str(chart_path), fig
            except Exception as e:
                logger.warning(f"Could not load chart image: {e}")

        # Fallback: create placeholder chart
        fig, ax = plt.subplots(figsize=(10, 6))
        ax.text(0.5, 0.5, f"Chart for {slide_key}\n{instruction[:50]}...",
                ha='center', va='center', fontsize=12, wrap=True)
        ax.set_title(f"Chart: {slide_key}")

        temp_dir = Path(tempfile.mkdtemp())
        chart_path = temp_dir / f"chart_{slide_key}.png"
        fig.savefig(chart_path, dpi=300, bbox_inches='tight')
        return str(chart_path), fig

    except Exception as e:
        logger.error(f"Error generating chart for {slide_key}: {str(e)}")
        return "", None


def generate_summary(
    summarizer_agent: SummarizerAgent, query: str, slide_key: str, instruction: str
) -> str:
    """Generate summary using summarizer agent.

    Args:
        summarizer_agent: Compiled summarizer agent
        query: Formatted query for summary generation
        slide_key: Unique identifier for the slide
        instruction: Summary instruction for fallback

    Returns:
        Generated summary text
    """
    try:
        summary_result = asyncio.run(summarizer_agent.invoke(query))
        summary_text = extract_summary_text(summary_result)
        return summary_text if summary_text else create_fallback_summary(slide_key, instruction)
    except Exception as e:
        logger.error(f"Error generating summary for {slide_key}: {str(e)}")
        return create_fallback_summary(slide_key, instruction)


def run_qa_review(
    critic_agent: CriticAgent, user_prompt_template: Any, qa_params: dict[str, Any],
    slide_title: str, chart_path: str | None, summary_text: str, config: dict[str, Any]
) -> dict[str, Any]:
    """Run QA review using critic agent.

    Args:
        critic_agent: Compiled critic agent
        user_prompt_template: Template for QA prompt
        qa_params: Quality assurance parameters
        slide_title: Title of the slide being reviewed
        chart_path: Path to generated chart
        summary_text: Generated summary text
        config: Slide configuration

    Returns:
        QA review result dictionary
    """
    try:
        chart_available = 'Available' if chart_path and Path(chart_path).exists() else 'Not available'
        summary_preview = summary_text[:300] if summary_text else 'Not available'

        slide_content = f"Slide Title: {slide_title}\nGenerated Chart: {chart_available}\nGenerated Summary: {summary_preview}"
        expected_requirements = f"Expected Slide Title: {slide_title}\nExpected Chart Instruction: {config.get('chart_instruction', '')}\nExpected Summary Instruction: {config.get('summary_instruction', '')}"

        qa_query = user_prompt_template.format(
            slide_content=slide_content,
            expected_requirements=expected_requirements,
            quality_standards=qa_params.get('quality_standards', ''),
            review_criteria=qa_params.get('review_criteria', '')
        )
        return asyncio.run(critic_agent.invoke(qa_query))
    except Exception as e:
        logger.error(f"Error in QA review: {str(e)}")
        return {"success": False, "error": str(e)}


def create_slide_presentation(
    title: str, chart_path: str | None, summary: str,
    layout_params: dict[str, Any], styling_params: dict[str, Any]
) -> Presentation:
    """Create a single slide presentation.

    Args:
        title: Slide title
        chart_path: Path to chart image
        summary: Summary text
        layout_params: Layout configuration
        styling_params: Styling configuration

    Returns:
        Presentation object with single slide
    """
    chart = chart_path if chart_path and Path(chart_path).exists() else ""
    return create_slide(
        slide_title=title, chart_path=chart, summary_text=summary,
        layout_params=layout_params, styling_params=styling_params
    )


def create_error_presentation(error_message: str) -> Presentation:
    """Create error presentation when workflow fails.

    Args:
        error_message: Error message to display

    Returns:
        Presentation object with error slide
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Multi-Agent Workflow Error"
    slide.placeholders[1].text = f"Error: {error_message}\nPlease check the logs for details."
    return prs
