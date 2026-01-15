"""Postprocessing functions for creating PowerPoint presentations.

This module contains deterministic functions for assembling presentations
from agent-generated content. These functions are reusable by both MA and SA pipelines.
"""
from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from ppt_autogen_workflow.base.utils import SYSTEM_FONT

logger = logging.getLogger(__name__)


def _get_blank_layout(prs: Presentation):
    """Find a blank slide layout from the presentation."""
    if len(prs.slide_layouts) > 6:
        return prs.slide_layouts[6]

    for layout in prs.slide_layouts:
        if "blank" in layout.name.lower():
            return layout

    return prs.slide_layouts[0]


def _add_title_box(
    slide,
    title_text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str,
    font_size: int,
    color_rgb: tuple[int, int, int],
) -> None:
    """Add a styled title to a slide."""
    title_box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    text_frame = title_box.text_frame
    text_frame.clear()

    p = text_frame.paragraphs[0]
    p.level = 0

    run = p.add_run()
    run.text = title_text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*color_rgb)


def _add_summary_box(
    slide,
    summary_text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str,
    font_size: int,
    paragraph_spacing: int,
    color_rgb: tuple[int, int, int],
) -> None:
    """Add formatted summary text box to slide."""
    text_box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_bottom = Inches(0.1)

    lines = summary_text.split("\n")
    text_frame.clear()

    para_count = 0
    for line in lines:
        line = line.strip()
        if not line:
            continue

        if para_count == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        cleaned_line = line[1:].strip() if line.startswith("•") else line

        run = p.add_run()
        run.text = f"• {cleaned_line}"
        run.font.size = Pt(font_size)
        run.font.name = font_name
        run.font.color.rgb = RGBColor(*color_rgb)

        p.level = 0
        p.space_after = Pt(paragraph_spacing)

        para_count += 1


def _copy_shape_to_slide(shape, new_slide) -> None:
    """Copy a shape from source slide to new slide."""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_stream = shape.image.blob
            new_slide.shapes.add_picture(
                io.BytesIO(image_stream),
                shape.left, shape.top, shape.width, shape.height,
            )

        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or (
            hasattr(shape, "text_frame") and shape.has_text_frame
        ):
            text_box = new_slide.shapes.add_textbox(
                shape.left, shape.top, shape.width, shape.height
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.clear()

            for i, src_paragraph in enumerate(shape.text_frame.paragraphs):
                if i == 0:
                    dst_paragraph = text_frame.paragraphs[0]
                else:
                    dst_paragraph = text_frame.add_paragraph()

                dst_paragraph.level = src_paragraph.level

                for src_run in src_paragraph.runs:
                    dst_run = dst_paragraph.add_run()
                    dst_run.text = src_run.text

                    try:
                        if src_run.font.size:
                            dst_run.font.size = src_run.font.size
                        if src_run.font.name:
                            dst_run.font.name = src_run.font.name
                        if src_run.font.bold is not None:
                            dst_run.font.bold = src_run.font.bold
                        if src_run.font.color and src_run.font.color.type is not None:
                            try:
                                if src_run.font.color.rgb:
                                    dst_run.font.color.rgb = src_run.font.color.rgb
                            except (AttributeError, TypeError):
                                pass
                    except (AttributeError, TypeError):
                        pass

                if not src_paragraph.runs:
                    dst_run = dst_paragraph.add_run()
                    dst_run.text = src_paragraph.text

    except Exception as e:
        logger.debug(f"Failed to copy shape: {e}")


def create_slide(
    slide_title: str,
    chart_path: str | Path,
    summary_text: str,
    layout_params: dict[str, Any] | None = None,
    styling_params: dict[str, Any] | None = None,
) -> Presentation:
    """Create a PowerPoint slide with chart and summary.

    This is a pure function that returns a Presentation object.

    Args:
        slide_title: Title for the slide
        chart_path: Path to the chart image file
        summary_text: Summary text content (bullet points)
        layout_params: Layout parameters (positions, dimensions)
        styling_params: Styling parameters (fonts, colors)

    Returns:
        python-pptx Presentation object containing the slide
    """
    layout = layout_params or {}
    styling = styling_params or {}

    # Layout defaults
    content_left = layout.get("content_left", 0.5)
    content_top = layout.get("content_top", 1.3)
    chart_width = layout.get("chart_width", 5.0)
    chart_height = layout.get("chart_height", 5.0)
    summary_width = layout.get("summary_width", 4.0)
    summary_spacing = layout.get("summary_spacing", 0.5)
    title_left = layout.get("title_left", 0.5)
    title_top = layout.get("title_top", 0.3)
    title_width = layout.get("title_width", 9.0)
    title_height = layout.get("title_height", 0.8)

    # Styling defaults
    title_font = styling.get("title_font", SYSTEM_FONT)
    title_size = styling.get("title_size", 28)
    text_font = styling.get("text_font", SYSTEM_FONT)
    text_size = styling.get("text_size_small", 11)
    paragraph_spacing = styling.get("paragraph_spacing", 8)
    title_color = styling.get("title_color_rgb", (31, 78, 121))
    text_color = styling.get("text_color_rgb", (51, 51, 51))

    # Create presentation
    prs = Presentation()
    slide_layout = _get_blank_layout(prs)
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    _add_title_box(
        slide,
        slide_title,
        title_left, title_top, title_width, title_height,
        title_font, title_size, title_color,
    )

    # Add chart (if path is provided and file exists)
    chart_added = False
    if chart_path:
        chart_path = Path(chart_path)
        if chart_path.exists():
            slide.shapes.add_picture(
                str(chart_path),
                Inches(content_left),
                Inches(content_top),
                Inches(chart_width),
                Inches(chart_height),
            )
            chart_added = True
        else:
            logger.warning(f"Chart path does not exist: {chart_path}")

    # Add summary
    if summary_text:
        if chart_added:
            # Summary on the right side when chart is present
            summary_left = content_left + chart_width + summary_spacing
            summary_w = summary_width
        else:
            # Summary takes more space when no chart
            summary_left = content_left
            summary_w = chart_width + summary_width + summary_spacing

        _add_summary_box(
            slide,
            summary_text,
            summary_left, content_top, summary_w, chart_height,
            text_font, text_size, paragraph_spacing, text_color,
        )

    logger.info(f"Created slide: {slide_title}")
    return prs


def combine_presentations(
    presentations: list[Presentation],
) -> Presentation:
    """Combine multiple presentations into one.

    This is a pure function that returns a combined Presentation object.

    Args:
        presentations: List of Presentation objects to combine

    Returns:
        Combined Presentation object
    """
    combined_prs = Presentation()
    slides_combined = 0

    for source_prs in presentations:
        for source_slide in source_prs.slides:
            slide_layout = _get_blank_layout(combined_prs)
            new_slide = combined_prs.slides.add_slide(slide_layout)

            # Copy all shapes
            for shape in source_slide.shapes:
                _copy_shape_to_slide(shape, new_slide)

            slides_combined += 1

    logger.info(f"Combined {slides_combined} slides into final presentation")
    return combined_prs


def assemble_presentation(
    slide_content: dict[str, Any],
    layout_params: dict[str, Any] | None = None,
    styling_params: dict[str, Any] | None = None,
) -> Any:
    """Assemble final presentation from generated slide content.

    This is a deterministic function that creates slides from the agent-generated content.
    Works for both MA and SA pipelines by handling different input structures.

    Args:
        slide_content: Dict containing slides. Can be either:
            - Direct dict mapping slide_key to content (SA format)
            - Dict with 'slides' key containing the slide dict (MA format)
        layout_params: Optional layout parameters for presentation
        styling_params: Optional styling parameters for presentation

    Returns:
        Combined PowerPoint presentation
    """
    try:
        # Handle both formats: MA wraps in 'slides' key, SA is direct dict
        slides = slide_content.get('slides', slide_content)

        slide_presentations = []
        for slide_key, content in slides.items():
            slide_title = content['slide_title']
            chart_path = content.get('chart_path', '')
            summary = content.get('summary', '')

            chart = chart_path if chart_path and Path(chart_path).exists() else ""
            slide_prs = create_slide(
                slide_title=slide_title,
                chart_path=chart,
                summary_text=summary,
                layout_params=layout_params,
                styling_params=styling_params,
            )
            slide_presentations.append(slide_prs)

        return combine_presentations(slide_presentations)

    except Exception as e:
        logger.error(f"Presentation assembly failed: {str(e)}", exc_info=True)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Presentation Assembly Error"
        slide.placeholders[1].text = f"Error: {str(e)}"
        return prs
